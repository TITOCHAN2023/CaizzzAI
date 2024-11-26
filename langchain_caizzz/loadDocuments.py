import os
import pandas as pd
import pdfplumber
from docx import Document as DocxDocument
from bs4 import BeautifulSoup
from pptx import Presentation
import pytesseract
from PIL import Image


def extract_text_from_file(filepath):
    """根据文件类型提取文本"""
    ext = os.path.splitext(filepath)[-1].lower()
    if ext == ".txt":
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
        
    elif ext == ".pdf":
        with pdfplumber.open(filepath) as pdf:
            return "\n".join(page.extract_text() for page in pdf.pages)
        
    elif ext in [".docx",".doc"]:
        doc = DocxDocument(filepath)
        return "\n".join(para.text for para in doc.paragraphs)
    
    elif ext == ".html":
        with open(filepath, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
            return soup.get_text()
        
    elif ext in [".xlsx", ".xls"]:
        df = pd.read_excel(filepath)
        return df.to_string(index=False)
    
    elif ext == ".pptx":
        prs = Presentation(filepath)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    
    elif ext in [".png", ".jpg", ".jpeg"]:
        text = pytesseract.image_to_string(Image.open(filepath))
        
        return text
    else:
        print(f"Unsupported file type: {ext}")
        return None
